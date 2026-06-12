"""
export.py — one-click leadership .pptx deck (the analyst's actual deliverable).

build_pptx(kpis, memo, gap_table, alerts, figures) -> bytes.
Slides: Title -> KPIs -> Action Center (memo + recruitment table + alerts) ->
one slide per key chart. Charts are rendered to PNG via kaleido. GUARDRAIL: if
image export fails, we fall back to a text-only deck so the download NEVER errors.
"""

import io
import threading

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from insights import inr

BLUE = RGBColor(0x25, 0x63, 0xEB)
DARK = RGBColor(0x0F, 0x17, 0x2A)
GREY = RGBColor(0x47, 0x55, 0x69)
RED = RGBColor(0xDC, 0x26, 0x26)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)


def _render_png(fig, width=1100, height=620, scale=2, timeout=8):
    """Render a plotly fig to PNG in a daemon thread with a hard timeout.

    kaleido can HANG (not error) on some machines; a plain try/except wouldn't
    save us. We join the render thread with a timeout and return None if it
    doesn't finish — so the deck build can never freeze the app live.
    """
    result = {}

    def _work():
        try:
            result["png"] = fig.to_image(format="png", width=width, height=height, scale=scale)
        except Exception:
            result["png"] = None

    t = threading.Thread(target=_work, daemon=True)
    t.start()
    t.join(timeout)
    return result.get("png")


# cache the one-time probe so we don't pay the timeout repeatedly
_KALEIDO_OK = None


def _kaleido_works():
    """Probe once with a tiny figure; if it returns a PNG fast, kaleido is usable."""
    global _KALEIDO_OK
    if _KALEIDO_OK is not None:
        return _KALEIDO_OK
    try:
        import plotly.graph_objects as go
        probe = go.Figure(go.Bar(x=[1], y=[1]))
        _KALEIDO_OK = _render_png(probe, width=200, height=120, scale=1, timeout=6) is not None
    except Exception:
        _KALEIDO_OK = False
    return _KALEIDO_OK


def _try_png(fig):
    if not _kaleido_works():
        return None
    return _render_png(fig)


def _title_slide(prs, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(12), Inches(2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "BD Performance & Action Report"
    r.font.size = Pt(40)
    r.font.bold = True
    r.font.color.rgb = DARK
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = subtitle
    r2.font.size = Pt(18)
    r2.font.color.rgb = BLUE
    return slide


def _kpi_slide(prs, kpis):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _heading(slide, "Key Performance Indicators")
    items = [
        ("GMV", inr(kpis["gmv"])),
        ("Bookings", f"{kpis['total_bookings']:,}"),
        ("AOV", inr(kpis["aov"])),
        ("Active partners", str(kpis["active_partners"])),
        ("Completion rate", f"{kpis['completion_rate']*100:.0f}%"),
        ("₹/week at risk", inr(kpis["revenue_at_risk_per_week"])),
    ]
    x0, y0, w, h = 0.7, 1.6, 4.0, 1.5
    for i, (label, val) in enumerate(items):
        col = i % 3
        row = i // 3
        box = slide.shapes.add_textbox(Inches(x0 + col * 4.1), Inches(y0 + row * 2.0),
                                       Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = val
        r.font.size = Pt(30); r.font.bold = True
        r.font.color.rgb = RED if label == "₹/week at risk" else BLUE
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = label
        r2.font.size = Pt(14); r2.font.color.rgb = GREY
    return slide


def _action_slide(prs, memo, gap_table, alerts):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _heading(slide, "Action Center — Where to recruit & ₹ at risk")

    # recruitment table (top undersupplied cells)
    under = None
    if gap_table is not None and len(gap_table):
        under = gap_table[gap_table["gap_flag"] == "UNDERSUPPLIED"].head(6)

    top = 1.4
    if under is not None and len(under):
        rows = len(under) + 1
        tbl = slide.shapes.add_table(rows, 4, Inches(0.6), Inches(top),
                                     Inches(8.0), Inches(0.4 * rows)).table
        headers = ["City", "Category", "Recruit", "₹/wk at risk"]
        for j, htext in enumerate(headers):
            c = tbl.cell(0, j); c.text = htext
            c.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
            c.text_frame.paragraphs[0].runs[0].font.bold = True
            c.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255)
            c.fill.solid(); c.fill.fore_color.rgb = BLUE
        for i, (_, r) in enumerate(under.iterrows(), start=1):
            vals = [r["city"], r["service_category"],
                    str(int(r["recommended_partner_adds"])),
                    inr(r["est_revenue_at_risk_per_week"])]
            for j, v in enumerate(vals):
                c = tbl.cell(i, j); c.text = str(v)
                c.text_frame.paragraphs[0].runs[0].font.size = Pt(11)

    # alerts on the right
    abox = slide.shapes.add_textbox(Inches(8.9), Inches(top), Inches(4.0), Inches(5))
    tf = abox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = "Alerts"
    r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = DARK
    for a in (alerts or [])[:8]:
        pp = tf.add_paragraph(); rr = pp.add_run()
        icon = "🔴 " if a.get("severity") == "high" else "🟠 "
        rr.text = icon + a["text"]
        rr.font.size = Pt(11)
        rr.font.color.rgb = RED if a.get("severity") == "high" else AMBER
    return slide


def _memo_slide(prs, memo):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _heading(slide, "Leadership Memo")
    box = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12.2), Inches(6))
    tf = box.text_frame; tf.word_wrap = True
    first = True
    for line in memo.splitlines():
        line = line.replace("###", "").replace("**", "").strip()
        if not line:
            continue
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run(); run.text = line
        is_head = line.endswith(":") or line.startswith("Top ") or line.startswith("Recommended") or line.startswith("Total")
        run.font.size = Pt(14 if is_head else 12)
        run.font.bold = is_head
        run.font.color.rgb = DARK if is_head else GREY
    return slide


def _chart_slide(prs, title, png):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _heading(slide, title)
    if png:
        slide.shapes.add_picture(io.BytesIO(png), Inches(0.8), Inches(1.4), width=Inches(11.5))
    else:
        box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1))
        box.text_frame.text = "(chart image unavailable — see live dashboard)"
    return slide


def _heading(slide, text):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.8))
    p = box.text_frame.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = DARK


def build_pptx(kpis, memo, gap_table, alerts, figures, subtitle="Automated Monday report"):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _title_slide(prs, subtitle)
    _kpi_slide(prs, kpis)
    _action_slide(prs, memo, gap_table, alerts)
    _memo_slide(prs, memo)

    chart_titles = {
        "demand_supply": "Demand vs Supply — recruit where supply is short",
        "revenue_by_city": "Revenue by city",
        "revenue_by_category": "Revenue share by category",
        "wow_trend": "Week-over-week trend + forecast",
        "onboarding_funnel": "Partner onboarding funnel",
        "status_by_category": "Booking outcomes by category",
    }
    # Only attempt chart slides if kaleido is responsive (probed once). Otherwise
    # ship the clean text+table deck — the live dashboard carries the visuals.
    if figures and _kaleido_works():
        for key, title in chart_titles.items():
            fig = figures.get(key)
            if fig is None:
                continue
            png = _try_png(fig)
            if png:
                _chart_slide(prs, title, png)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def build_pdf(kpis, memo, gap_table, alerts):
    """Optional text-only PDF alternative (no chart images needed)."""
    try:
        from fpdf import FPDF
        from fpdf.enums import XPos, YPos
    except Exception:
        return None
    try:
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()
        w = pdf.epw  # effective page width (page minus margins)

        pdf.set_font("Helvetica", "B", 18)
        pdf.multi_cell(w, 12, "BD Performance & Action Report",
                       new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.set_font("Helvetica", "", 11)
        pdf.multi_cell(w, 8,
                       f"GMV {_ascii(inr(kpis['gmv']))} | Bookings {kpis['total_bookings']:,} | "
                       f"At risk/wk {_ascii(inr(kpis['revenue_at_risk_per_week']))}",
                       new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.ln(3)

        pdf.set_font("Helvetica", "", 10)
        for line in memo.splitlines():
            line = _ascii(line.replace("###", "").replace("**", "").strip())
            if line:
                pdf.multi_cell(w, 6, line, new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        return bytes(pdf.output())
    except Exception:
        return None


def _ascii(s):
    return (s.replace("₹", "Rs.").replace("≈", "~").replace("·", "-")
            .encode("latin-1", "ignore").decode("latin-1"))
